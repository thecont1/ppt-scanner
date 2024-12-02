import os
from pptx import Presentation
import pandas as pd
import sys
import openai
import time
from dotenv import load_dotenv
from tqdm import tqdm

# Load environment variables from .env file
load_dotenv()

def extract_slide_text(slide):
    """Extract text from a slide and check for images"""
    text = []
    has_image = False
    
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text.append(shape.text.strip())
        # Check if shape is an image
        if shape.shape_type in [13, 14, 17, 19]:  # Common image shape types
            has_image = True
    
    text = ' '.join(text)
    if not text and has_image:
        return "[Image Slide]"
    return text

def get_key_phrase(text, api_key):
    """Extract a key phrase (up to 12 words) from the slide content"""
    if not text.strip() or text == "[Image Slide]":
        return text
    
    try:
        openai.api_key = api_key
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that extracts key phrases from text. Select a consecutive string of words (maximum 12 words) from the input text that best represents its main topic. Do not generate new text - only extract existing words in the same order they appear."},
                {"role": "user", "content": f"Extract a key phrase (max 12 consecutive words) from this text that hints at its contents: {text}"}
            ],
            max_tokens=50,
            temperature=0
        )
        return response.choices[0].message.content
    except Exception as e:
        if "insufficient_quota" in str(e) or "billing" in str(e).lower():
            print("\nERROR: OpenAI API credit balance is too low or there's a billing issue.")
            print("Please check your API key and billing status.")
            sys.exit(1)  # Terminate program immediately
        raise  # Re-raise other exceptions

def get_ppt_info(directory, api_key):
    """
    Scan directory for PowerPoint files and collect information about each slide.
    """
    if not api_key:
        print("Error: OPENAI_API_KEY environment variable not set")
        sys.exit(1)
    
    ppt_data = []
    
    # Supported PowerPoint extensions
    ppt_extensions = ('.pptx', '.ppt')
    
    # Get list of PowerPoint files
    ppt_files = [f for f in os.listdir(directory) if f.lower().endswith(ppt_extensions)]
    
    if not ppt_files:
        print("No PowerPoint files found in the specified directory!")
        return []
    
    print(f"Found {len(ppt_files)} PowerPoint file(s)")
    
    # Process each PowerPoint file with progress bar
    for filename in tqdm(ppt_files, desc="Processing files", unit="file"):
        filepath = os.path.join(directory, filename)
        
        try:
            # Process each slide in the presentation
            prs = Presentation(filepath)
            slides = list(prs.slides)
            
            # Inner progress bar for slides
            for slide_number, slide in enumerate(tqdm(slides, desc=f"Processing {filename}", unit="slide", leave=False), 1):
                # Extract text from the slide
                slide_text = extract_slide_text(slide)
                
                # Get key phrase using OpenAI API
                key_phrase = get_key_phrase(slide_text, api_key)
                
                # Add data to list
                ppt_data.append({
                    'Filename': filename,
                    'Slide Number': slide_number,
                    'Key Phrase': key_phrase
                })
                
                # Add a small delay to avoid hitting API rate limits
                time.sleep(0.5)
                
        except Exception as e:
            print(f"Error processing {filename}: {str(e)}")
            if "insufficient_quota" in str(e) or "billing" in str(e).lower():
                sys.exit(1)  # Terminate on API quota issues
    
    return ppt_data

def main():
    # Get directory path from command line argument or use current directory
    if len(sys.argv) > 1:
        directory = sys.argv[1]
    else:
        directory = os.path.dirname(os.path.abspath(__file__))
    
    # Get OpenAI API key from environment variable
    api_key = os.getenv('OPENAI_API_KEY')
    if not api_key:
        print("Error: OPENAI_API_KEY environment variable not set")
        return
    
    if not os.path.exists(directory):
        print(f"Directory does not exist: {directory}")
        return
    
    print(f"Scanning directory: {directory}")
    
    # Get PowerPoint files information
    ppt_data = get_ppt_info(directory, api_key)
    
    if not ppt_data:
        return
    
    print("\nGenerating Excel file...")
    
    # Create DataFrame and export to Excel
    df = pd.DataFrame(ppt_data)
    output_file = os.path.join(directory, 'ppt_slides.xlsx')
    
    # Create Excel writer with xlsxwriter engine
    with pd.ExcelWriter(output_file, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='Sheet1')
        
        # Get the xlsxwriter workbook and worksheet objects
        workbook = writer.book
        worksheet = writer.sheets['Sheet1']
        
        # Set column widths
        worksheet.set_column('A:A', 30)  # Filename
        worksheet.set_column('B:B', 15)  # Slide Number
        worksheet.set_column('C:C', 30)  # Key Phrase
    
    print(f"\nAnalysis complete! Results saved to: {output_file}")

if __name__ == "__main__":
    main()
